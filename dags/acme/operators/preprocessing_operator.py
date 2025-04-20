from typing import Any, Dict, Union, List, Optional
from abc import ABC, abstractmethod
from airflow.models import BaseOperator
from airflow.utils.decorators import apply_defaults
from PIL import Image
import pandas as pd
import io
import pytesseract
from pdf2image import convert_from_path
import mammoth
import openpyxl
import csv
from pptx import Presentation
import os
from pathlib import Path
import numpy as np
import cv2

class PreprocessingStrategy(ABC):
    """Abstract base class for preprocessing strategies"""
    
    @abstractmethod
    def preprocess(self, input_data: Any) -> Dict[str, Any]:
        """
        Preprocess the input data and return a dictionary containing:
        - processed_data: The preprocessed data
        - metadata: Any additional information about the processing
        """
        pass

class ImagePreprocessingStrategy(PreprocessingStrategy):
    """Strategy for preprocessing image inputs"""
    
    def __init__(self):
        self.supported_formats = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.svg'}
        
    def _resize_image(self, image: Image.Image, max_size: int = 1024) -> Image.Image:
        """Resize image while maintaining aspect ratio"""
        ratio = max_size / max(image.size)
        if ratio < 1:
            new_size = tuple(int(dim * ratio) for dim in image.size)
            return image.resize(new_size, Image.Resampling.LANCZOS)
        return image

    def _enhance_image(self, image: np.ndarray) -> np.ndarray:
        """Enhance image quality"""
        # Convert to grayscale
        if len(image.shape) == 3:
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        else:
            gray = image
            
        # Apply adaptive histogram equalization
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
        enhanced = clahe.apply(gray)
        
        # Denoise
        denoised = cv2.fastNlMeansDenoising(enhanced)
        
        return denoised

    def preprocess(self, input_data: Union[str, bytes, Image.Image]) -> Dict[str, Any]:
        # Convert input to PIL Image if needed
        if isinstance(input_data, str):
            if not any(input_data.lower().endswith(ext) for ext in self.supported_formats):
                raise ValueError(f"Unsupported image format. Supported formats: {self.supported_formats}")
            image = Image.open(input_data)
        elif isinstance(input_data, bytes):
            image = Image.open(io.BytesIO(input_data))
        elif isinstance(input_data, Image.Image):
            image = input_data
        else:
            raise ValueError("Unsupported image input type")

        # Preprocess image
        image = self._resize_image(image)
        
        # Convert to numpy array for OpenCV processing
        np_image = np.array(image)
        enhanced_image = self._enhance_image(np_image)
        
        # Extract text using OCR if applicable
        try:
            extracted_text = pytesseract.image_to_string(enhanced_image)
        except:
            extracted_text = ""

        # Convert back to PIL Image
        processed_image = Image.fromarray(enhanced_image)

        return {
            'processed_data': {
                'image': processed_image,
                'extracted_text': extracted_text
            },
            'metadata': {
                'original_size': image.size,
                'processed_size': processed_image.size,
                'format': image.format,
                'has_text': bool(extracted_text.strip())
            }
        }

class OfficeDocumentPreprocessingStrategy(PreprocessingStrategy):
    """Strategy for preprocessing Microsoft Office and similar documents"""
    
    def __init__(self):
        self.supported_formats = {
            '.docx': self._process_docx,
            '.xlsx': self._process_xlsx,
            '.pptx': self._process_pptx,
            '.pdf': self._process_pdf,
            '.csv': self._process_csv
        }

    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        with open(file_path, 'rb') as docx_file:
            result = mammoth.convert_to_markdown(docx_file)
            return {
                'text': result.value,
                'messages': result.messages
            }

    def _process_xlsx(self, file_path: str) -> Dict[str, Any]:
        wb = openpyxl.load_workbook(file_path)
        sheets_data = {}
        
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            data = []
            for row in ws.rows:
                data.append([str(cell.value) if cell.value is not None else '' for cell in row])
            sheets_data[sheet] = pd.DataFrame(data[1:], columns=data[0])
            
        return {
            'sheets': sheets_data
        }

    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        prs = Presentation(file_path)
        slides_data = []
        
        for slide in prs.slides:
            slide_content = {
                'text': [],
                'images': []
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content['text'].append(shape.text)
                if hasattr(shape, "image"):
                    # Save image to temporary file
                    image_stream = io.BytesIO(shape.image.blob)
                    image = Image.open(image_stream)
                    slide_content['images'].append(image)
                    
            slides_data.append(slide_content)
            
        return {
            'slides': slides_data
        }

    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        # Convert PDF to images
        images = convert_from_path(file_path)
        
        pages_data = []
        for image in images:
            # Extract text using OCR
            text = pytesseract.image_to_string(image)
            pages_data.append({
                'image': image,
                'text': text
            })
            
        return {
            'pages': pages_data
        }

    def _process_csv(self, file_path: str) -> Dict[str, Any]:
        df = pd.read_csv(file_path)
        return {
            'dataframe': df
        }

    def preprocess(self, input_data: str) -> Dict[str, Any]:
        file_ext = Path(input_data).suffix.lower()
        
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported document format. Supported formats: {self.supported_formats.keys()}")
            
        processor = self.supported_formats[file_ext]
        processed_data = processor(input_data)
        
        return {
            'processed_data': processed_data,
            'metadata': {
                'file_type': file_ext,
                'processor_used': processor.__name__
            }
        }

class RDBPreprocessingStrategy(PreprocessingStrategy):
    """Strategy for preprocessing RDB data"""
    
    def _clean_column_names(self, df: pd.DataFrame) -> pd.DataFrame:
        """Clean and standardize column names"""
        df.columns = df.columns.str.lower()
        df.columns = df.columns.str.replace('[^a-z0-9_]', '_', regex=True)
        return df

    def _handle_missing_values(self, df: pd.DataFrame) -> pd.DataFrame:
        """Handle missing values intelligently"""
        # For numeric columns, fill with median
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        df[numeric_cols] = df[numeric_cols].fillna(df[numeric_cols].median())
        
        # For categorical columns, fill with mode
        categorical_cols = df.select_dtypes(include=['object']).columns
        for col in categorical_cols:
            df[col] = df[col].fillna(df[col].mode()[0] if not df[col].mode().empty else 'UNKNOWN')
            
        return df

    def _detect_and_handle_outliers(self, df: pd.DataFrame) -> pd.DataFrame:
        """Detect and handle outliers in numeric columns"""
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        
        for col in numeric_cols:
            Q1 = df[col].quantile(0.25)
            Q3 = df[col].quantile(0.75)
            IQR = Q3 - Q1
            lower_bound = Q1 - 1.5 * IQR
            upper_bound = Q3 + 1.5 * IQR
            
            # Cap outliers at boundaries
            df[col] = df[col].clip(lower=lower_bound, upper=upper_bound)
            
        return df

    def preprocess(self, input_data: Union[pd.DataFrame, Dict, str]) -> Dict[str, Any]:
        if isinstance(input_data, str):
            # Assume it's a SQL query result
            try:
                df = pd.read_sql(input_data, connection)  # connection should be provided
            except:
                raise ValueError("Invalid SQL query or connection not available")
        elif isinstance(input_data, dict):
            df = pd.DataFrame(input_data)
        elif isinstance(input_data, pd.DataFrame):
            df = input_data.copy()
        else:
            raise ValueError("Unsupported RDB data input type")

        # Apply preprocessing steps
        df = self._clean_column_names(df)
        df = self._handle_missing_values(df)
        df = self._detect_and_handle_outliers(df)

        # Generate statistics
        statistics = {
            'row_count': len(df),
            'column_count': len(df.columns),
            'missing_values': df.isnull().sum().to_dict(),
            'data_types': df.dtypes.astype(str).to_dict()
        }

        return {
            'processed_data': df,
            'metadata': {
                'statistics': statistics,
                'preprocessing_steps': [
                    'column_name_cleaning',
                    'missing_value_handling',
                    'outlier_detection_and_handling'
                ]
            }
        }

class PreprocessingOperator(BaseOperator):
    """
    Custom Airflow operator that handles preprocessing of various input types
    
    :param input_data: Input data to be preprocessed
    :param input_type: Type of input ('image', 'office_document', or 'rdb')
    :param preprocessing_params: Additional parameters for preprocessing
    """
    
    @apply_defaults
    def __init__(
        self,
        input_data: Any,
        input_type: str,
        preprocessing_params: Optional[Dict] = None,
        **kwargs
    ) -> None:
        super().__init__(**kwargs)
        self.input_data = input_data
        self.input_type = input_type
        self.preprocessing_params = preprocessing_params or {}
        self._initialize_strategies()

    def _initialize_strategies(self) -> None:
        """Initialize preprocessing strategies"""
        self.strategies = {
            'image': ImagePreprocessingStrategy(),
            'office_document': OfficeDocumentPreprocessingStrategy(),
            'rdb': RDBPreprocessingStrategy()
        }

    def execute(self, context: Dict) -> Dict[str, Any]:
        """
        Execute the preprocessing operation
        
        :param context: Airflow context
        :return: Dictionary containing processed data and metadata
        """
        # Pull input data from XCom
        self.input_data = context['task_instance'].xcom_pull(task_ids=self.input_data)
        
        # Initialize results dictionary
        results = {}

        # Check if input_data is a dictionary
        if isinstance(self.input_data, dict):
            for input_type, files in self.input_data.items():
                if input_type not in self.strategies:
                    self.log.warning(f"Unsupported input type: {input_type}")
                    continue
                
                strategy = self.strategies[input_type]
                
                for file_path in files:
                    try:
                        self.log.info(f"Preprocessing {input_type} input: {file_path}")
                        result = strategy.preprocess(file_path)
                        results[file_path] = result
                        self.log.info(f"Preprocessing completed successfully for {file_path}")
                    except Exception as e:
                        self.log.error(f"Preprocessing failed for {file_path}: {str(e)}")
        else:
            raise ValueError("input_data should be a dictionary with input types as keys")

        return results
